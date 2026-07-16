#!/usr/bin/env python3
"""Build product-overview.pptx as a slide-deck mirror of product-overview.html.

Content (headings, copy, bullets, screenshots) is kept in sync by hand with
product-overview.html — re-run this after editing that file's copy.
Requires: pip install python-pptx pillow

Usage: python3 build-pptx.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(HERE, 'assets', 'overview')
OUT = os.path.join(HERE, 'product-overview.pptx')

SERIF = 'Georgia'
SANS = 'Arial'

NAVY_DEEP = RGBColor.from_string('141D4A')
NAVY = RGBColor.from_string('1E2761')
GOLD = RGBColor.from_string('C9A84C')
GOLD_LT = RGBColor.from_string('E8D09A')
CREAM = RGBColor.from_string('F7F5EF')
WARM = RGBColor.from_string('F5F0E8')
INK = RGBColor.from_string('1A1A1A')
MUTED = RGBColor.from_string('6B6B6B')
WHITE = RGBColor.from_string('FFFFFF')
LINE = RGBColor.from_string('E3E1DA')

GREEN_BG, GREEN_TX = RGBColor.from_string('EAF3DE'), RGBColor.from_string('27500A')
RED_BG, RED_TX = RGBColor.from_string('FCEBEB'), RGBColor.from_string('791F1F')
AMBER_BG, AMBER_TX = RGBColor.from_string('FDF0DD'), RGBColor.from_string('8A5A12')
BLUE_BG, BLUE_TX = RGBColor.from_string('DDEAF8'), RGBColor.from_string('1E4D8C')
PURPLE_BG, PURPLE_TX = RGBColor.from_string('EEEDFE'), RGBColor.from_string('4A3280')
WHITE_08 = RGBColor.from_string('2A3268')  # translucent-ish white-on-navy panel approximation

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide(bg=CREAM):
    slide = prs.slides.add_slide(BLANK)
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    rect.fill.solid()
    rect.fill.fore_color.rgb = bg
    rect.line.fill.background()
    rect.shadow.inherit = False
    return slide


def no_autofit(tf):
    el = tf._txBody
    bodyPr = el.find(qn('a:bodyPr'))
    for tag in ('a:normAutofit', 'a:spAutoFit'):
        existing = bodyPr.find(qn(tag))
        if existing is not None:
            bodyPr.remove(existing)
    bodyPr.append(bodyPr.makeelement(qn('a:noAutofit'), {}))


def add_text(slide, left, top, width, height, text, size=14, color=INK, bold=False,
             font=SANS, align=PP_ALIGN.LEFT, italic=False, line_spacing=None, anchor=MSO_ANCHOR.TOP,
             letter_spacing=None, upper=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    no_autofit(tf)
    tf.vertical_anchor = anchor
    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line.upper() if upper else line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = font
        r.font.color.rgb = color
    return box


def add_rect(slide, left, top, width, height, fill=None, line_color=None, line_w=None, shadow=False, round_=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, left, top, width, height)
    if round_:
        try:
            shp.adjustments[0] = 0.06
        except Exception:
            pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = line_w or Pt(0.75)
    shp.shadow.inherit = shadow
    return shp


def add_pill(slide, left, top, text, fg, bg, size=9, width=None, height=Inches(0.28), bold=True):
    w = width or Inches(0.15 + 0.078 * len(text))
    shp = add_rect(slide, left, top, w, height, fill=bg, round_=True)
    try:
        shp.adjustments[0] = 0.5
    except Exception:
        pass
    tf = shp.text_frame
    tf.margin_left = tf.margin_right = Pt(4)
    tf.margin_top = tf.margin_bottom = Pt(1)
    no_autofit(tf)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text.upper()
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = SANS
    r.font.color.rgb = fg
    return shp


def add_bullets(slide, left, top, width, height, items, size=13, color=INK, gap=10, bullet_color=GOLD, bold_lead=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    no_autofit(tf)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.18
        r1 = p.add_run()
        r1.text = '✦  '
        r1.font.size = Pt(size)
        r1.font.color.rgb = bullet_color
        r1.font.name = SANS
        r2 = p.add_run()
        r2.text = item
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
        r2.font.name = SANS
    return box


def fit_image_box(path, max_w, max_h):
    with Image.open(path) as im:
        iw, ih = im.size
    ar = iw / ih
    box_ar = max_w / max_h
    if ar > box_ar:
        w = max_w
        h = Emu(int(max_w / ar))
    else:
        h = max_h
        w = Emu(int(max_h * ar))
    return w, h


def add_image_fit(slide, path, left, top, max_w, max_h, shadow=True, border=True, center_h=True, center_v=False):
    w, h = fit_image_box(path, max_w, max_h)
    l = left + (max_w - w) // 2 if center_h else left
    t = top + (max_h - h) // 2 if center_v else top
    pic = slide.shapes.add_picture(path, l, t, w, h)
    if border:
        pic.line.color.rgb = LINE
        pic.line.width = Pt(0.75)
    pic.shadow.inherit = False
    return pic, l, t, w, h


def eyebrow(slide, left, top, text, color=GOLD, size=11):
    add_text(slide, left, top, Inches(8), Inches(0.3), text, size=size, color=color, bold=True, upper=True)


def heading(slide, left, top, width, text, size=28, color=NAVY):
    add_text(slide, left, top, width, Inches(1.1), text, size=size, color=color, bold=False, font=SERIF)


MX = Inches(0.7)   # standard left/right margin
CW = SW - 2 * MX    # content width

# ─────────────────────────────────────────────────────────────────────────────
# 1. HERO
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide(NAVY_DEEP)
add_text(s, MX, Inches(1.55), CW, Inches(0.4), 'MAISON AURORE ROSTER', size=13, color=GOLD, bold=True,
         align=PP_ALIGN.CENTER, letter_spacing=True)
add_text(s, Inches(1.3), Inches(2.15), SW - Inches(2.6), Inches(2.0),
         'One system, four roles — the full store-operations\nloop from schedule generation to the floor.',
         size=32, color=WHITE, bold=False, font=SERIF, align=PP_ALIGN.CENTER, line_spacing=1.15)
add_text(s, Inches(2.3), Inches(3.85), SW - Inches(4.6), Inches(0.8),
         'Product overview and budget brief for stakeholders evaluating rollout across boutiques.',
         size=14, color=RGBColor.from_string('B9B8C9'), align=PP_ALIGN.CENTER)

stat_items = [('4', 'ROLE-BASED PORTALS'), ('8', 'CONFIGURABLE RULES'), ('2', 'BOUTIQUES IN THE CURRENT DATASET')]
stat_w = Inches(3.2)
total_w = stat_w * 3
start_x = (SW - total_w) // 2
for i, (n, l) in enumerate(stat_items):
    x = start_x + stat_w * i
    add_text(s, x, Inches(4.85), stat_w, Inches(0.75), n, size=34, color=GOLD, bold=False, font=SERIF, align=PP_ALIGN.CENTER)
    add_text(s, x, Inches(5.55), stat_w, Inches(0.5), l, size=10, color=RGBColor.from_string('8C8BA3'), bold=True,
              align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# 2. EXECUTIVE SUMMARY
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide(WHITE)
eyebrow(s, MX, Inches(0.45), 'EXECUTIVE SUMMARY')
heading(s, MX, Inches(0.78), CW, 'A scheduling system built around who actually touches it.', size=25)

copy_w = Inches(7.1)
add_bullets(s, MX, Inches(1.85), copy_w, Inches(5.2), [
    'Rebuilt from a single admin screen into three working, role-scoped portals — Admin, Approver and Staff — plus a fourth (Reader, a read-only published-roster browser) reserved for a later phase.',
    'Every roster moves through the same tracked lifecycle: draft → submitted → approved → published (with reject and amend paths), scored out of 100 by one scoring module shared by the generator and manual edits alike.',
    'Staff, VIC clients, shifts, rules and scoring weights are stored per boutique in one data model. The current dataset has two boutiques configured (Sydney and Melbourne); adding another is configuration, not a new build.',
], size=13.5)

card_x = MX + copy_w + Inches(0.35)
card_w = CW - copy_w - Inches(0.35)
cards = [
    ("WHO IT'S FOR", "Store or regional operations admins, the manager who signs off each roster, and every floor employee checking their own shifts."),
    ('WHAT CHANGED RECENTLY', 'Split into four dedicated portals with real sign-in; added submit → approve → publish, live-rescoring manual edits, leave tracking, and per-boutique rules.'),
    ('WHY IT MATTERS FOR BUDGET', 'One codebase and one data model serve every boutique — scaling to another store is configuration inside the existing schema, not a new build.'),
]
cy = Inches(1.85)
for k, v in cards:
    ch = Inches(1.55)
    bar = add_rect(s, card_x, cy, Inches(0.045), ch, fill=GOLD)
    box = add_rect(s, card_x + Inches(0.045), cy, card_w - Inches(0.045), ch, fill=WARM)
    add_text(s, card_x + Inches(0.2), cy + Inches(0.1), card_w - Inches(0.4), Inches(0.3), k, size=11, color=NAVY, bold=True)
    add_text(s, card_x + Inches(0.2), cy + Inches(0.42), card_w - Inches(0.4), ch - Inches(0.5), v, size=10.5, color=MUTED, line_spacing=1.15)
    cy += ch + Inches(0.22)

# ─────────────────────────────────────────────────────────────────────────────
# 3. PORTAL MAP
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide(WARM)
eyebrow(s, MX, Inches(0.45), "HOW IT'S ORGANIZED")
heading(s, MX, Inches(0.78), CW, 'Four portals, one login each.', size=25)
add_text(s, MX, Inches(1.5), CW, Inches(0.5),
         "Every user signs in once and lands in the portal that matches their role — no menu of features they don't need.",
         size=13, color=MUTED)

portals = [
    ('ADMIN', AMBER_BG, AMBER_TX, 'Operations control', 'Generate and edit rosters, manage staff, VIC clients, shifts, rules and leave — all scoped to the active boutique.', 'Store / regional admin'),
    ('APPROVER', AMBER_BG, AMBER_TX, 'Sign-off & publish', 'A focused inbox of rosters awaiting review, with approve, reject and publish — kept separate from whoever built the draft.', 'Boutique / regional manager'),
    ('STAFF', BLUE_BG, BLUE_TX, 'My Schedule', 'Employees see their own upcoming shifts straight from the published roster, no admin access required.', 'Floor staff'),
    ('READER', PURPLE_BG, PURPLE_TX, 'Published view (planned)', 'A read-only, PDF/CSV-exportable roster browser for a shared back-of-house screen or printout.', 'Shared / front-of-house'),
]
gap = Inches(0.25)
card_w = (CW - gap * 3) / 4
card_h = Inches(4.5)
cy = Inches(2.25)
for i, (tag, tbg, ttx, title, desc, who) in enumerate(portals):
    cx = MX + i * (card_w + gap)
    add_rect(s, cx, cy, card_w, card_h, fill=WHITE, line_color=LINE, line_w=Pt(0.75), round_=True)
    add_pill(s, cx + Inches(0.2), cy + Inches(0.22), tag, ttx, tbg)
    add_text(s, cx + Inches(0.2), cy + Inches(0.65), card_w - Inches(0.4), Inches(0.8), title, size=15, color=NAVY, bold=True, font=SANS, line_spacing=1.1)
    add_text(s, cx + Inches(0.2), cy + Inches(1.55), card_w - Inches(0.4), Inches(2.2), desc, size=10.5, color=MUTED, line_spacing=1.3)
    add_text(s, cx + Inches(0.2), cy + card_h - Inches(0.55), card_w - Inches(0.4), Inches(0.4), who, size=10, color=GOLD, bold=True)

# ─────────────────────────────────────────────────────────────────────────────
# 4. ROSTER LIFECYCLE (workflow diagram)
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide(WHITE)
eyebrow(s, MX, Inches(0.45), 'THE ROSTER LIFECYCLE')
heading(s, MX, Inches(0.78), CW, 'From a blank date to a shift someone sees on their phone.', size=22)
add_text(s, MX, Inches(1.5), CW, Inches(0.4),
          'Every roster — machine-generated or hand-built — moves through the same six steps, with a permanent record at each one.',
          size=12, color=MUTED)

steps = [
    ('1', 'ADMIN · ROSTERS', 'Generate a draft',
     "Pick a date. The engine reads that boutique's staff, skills, VIC appointments, leave and shift rules and proposes a scored roster."),
    ('2', 'ADMIN · ROSTERS', 'Review & adjust',
     'Unmet requirements and rule flags are called out per shift. Add or remove staff by hand — the score recalculates live.'),
    ('3', 'ADMIN → APPROVER', 'Submit for approval',
     "The draft leaves the admin's queue and appears in the approver's Approval Inbox, score and overrides visible."),
]
steps2 = [
    ('5', 'APPROVER · INBOX', 'Publish',
     "If another roster is already published for that boutique and date, it's automatically marked Amended and superseded."),
    ('6', 'STAFF · MY SCHEDULE', 'Live for staff',
     'Every linked staff member on that roster sees the shift on My Schedule the moment it publishes — no announcement needed.'),
]

col_w = (CW - Inches(0.4) * 2) / 3
top0 = Inches(2.1)
row_h = Inches(2.05)


def flow_card(x, y, w, h, num, role, title, body):
    circ = slide_shapes_add_oval(s, x, y, Inches(0.5), Inches(0.5))
    add_text(s, x, y, Inches(0.5), Inches(0.5), num, size=16, color=GOLD_LT, font=SERIF, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tx = x + Inches(0.65)
    tw = w - Inches(0.65)
    add_text(s, tx, y - Inches(0.02), tw, Inches(0.3), role, size=9, color=GOLD, bold=True)
    add_text(s, tx, y + Inches(0.28), tw, Inches(0.45), title, size=14.5, color=NAVY, bold=True)
    add_text(s, tx, y + Inches(0.72), tw, h - Inches(0.8), body, size=10.5, color=MUTED, line_spacing=1.22)


def slide_shapes_add_oval(slide, x, y, w, h, fill=NAVY):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


for i, (num, role, title, body) in enumerate(steps):
    x = MX + i * (col_w + Inches(0.4))
    flow_card(x, top0, col_w, row_h, num, role, title, body)

# Step 4 spans full width with the approve/reject branch
y4 = top0 + row_h + Inches(0.25)
flow_card(MX, y4, CW, Inches(1.0), '4', 'APPROVER · INBOX', 'Approve or reject',
          'The approver — a separate role from whoever built the draft — checks it and decides.')
branch_y = y4 + Inches(1.05)
branch_w = (CW - Inches(0.65) - Inches(0.3)) / 2
ok_box = add_rect(s, MX + Inches(0.65), branch_y, branch_w, Inches(0.62), fill=GREEN_BG, round_=True)
add_text(s, MX + Inches(0.85), branch_y + Inches(0.06), branch_w - Inches(0.4), Inches(0.5),
          'APPROVED\nMoves on to publish.', size=10, color=GREEN_TX, bold=False, line_spacing=1.1)
no_box = add_rect(s, MX + Inches(0.65) + branch_w + Inches(0.3), branch_y, branch_w, Inches(0.62), fill=RED_BG, round_=True)
add_text(s, MX + Inches(0.85) + branch_w + Inches(0.3), branch_y + Inches(0.06), branch_w - Inches(0.4), Inches(0.5),
          'REJECTED\nReturns to Draft with a note, back to step 2 — no dead end.', size=10, color=RED_TX, line_spacing=1.1)

y5 = branch_y + Inches(0.85)
for i, (num, role, title, body) in enumerate(steps2):
    x = MX + i * (col_w * 1.5 + Inches(0.4))
    flow_card(x, y5, col_w * 1.5, Inches(1.15), num, role, title, body)

# ─────────────────────────────────────────────────────────────────────────────
# Feature-slide helper
# ─────────────────────────────────────────────────────────────────────────────
IMG_MAX_W = Inches(5.55)
IMG_MAX_H = Inches(6.05)
IMG_TOP = Inches(0.9)


def feature_slide(bg, role_tag, title, body, bullets, image_path, image_right=True, caption=None, mobile=False):
    sl = add_slide(bg)
    text_w = CW - IMG_MAX_W - Inches(0.45)
    if image_right:
        img_x = SW - MX - IMG_MAX_W
        text_x = MX
    else:
        img_x = MX
        text_x = MX + IMG_MAX_W + Inches(0.45)

    img_max_h = Inches(5.4) if mobile else IMG_MAX_H
    pic, l, t, w, h = add_image_fit(sl, image_path, img_x, IMG_TOP, IMG_MAX_W, img_max_h, center_h=True)
    if caption:
        add_text(sl, img_x, t + h + Inches(0.1), IMG_MAX_W, Inches(0.5), caption, size=10, color=MUTED,
                  italic=True, align=PP_ALIGN.CENTER)

    add_text(sl, text_x, Inches(0.75), text_w, Inches(0.35), role_tag, size=10.5, color=GOLD, bold=True, upper=True)
    add_text(sl, text_x, Inches(1.12), text_w, Inches(1.35), title, size=21, color=NAVY, font=SERIF, line_spacing=1.15)
    add_text(sl, text_x, Inches(2.55), text_w, Inches(1.75), body, size=12, color=MUTED, line_spacing=1.32)
    add_bullets(sl, text_x, Inches(4.45), text_w, Inches(2.4), bullets, size=11.5, gap=9)
    return sl


A = os.path.join(ASSETS, '{}')

feature_slide(
    WHITE, 'ADMIN · ROSTERS',
    "Generate a full roster in one click, see exactly what's short.",
    'The engine reads active staff, skills, VIC appointments, leave and shift rules for the target date and proposes a complete roster — scored out of 100, with every unmet requirement surfaced immediately.',
    ['One-click generation per date, respecting leave and day-of-week availability',
     'Live skill / seniority / VIC-coverage / gender-balance breakdown per shift',
     'Draft → Submitted → Approved → Published status tracked with full history'],
    A.format('admin-roster-detail.png'), image_right=True,
    caption='Generated draft — scored automatically, with gaps called out per shift',
)

feature_slide(
    WHITE, 'ADMIN · ROSTERS',
    'Managers keep the final say, with the score watching in real time.',
    'Add or remove anyone from a shift and the score, skill coverage and VIC-coverage badges recompute instantly — using the same scoring logic the generator itself uses.',
    ['Add-from-bench and remove-in-place, no page reload',
     'Overrides are tracked and shown to the approver for context',
     'Nothing is submitted until a manager is satisfied with the result'],
    A.format('admin-manual-edit.png'), image_right=False,
    caption='Filling a gap by hand — the score updates before saving',
)

feature_slide(
    WHITE, 'ADMIN · STAFF',
    'One record per employee, shared by every downstream feature.',
    'Employment type, contracted hours, seniority, languages, primary skill and day-of-week availability all live on the staff record — the same data the scheduling engine, VIC assignment and leave tracker all read from.',
    ['Full-time, part-time, casual and contractor employment types',
     'Per-boutique day-of-week availability, not just a blanket flag',
     'Search by name, HR ID or skill'],
    A.format('admin-staff.png'), image_right=True,
    caption='Staff directory — employment type, skills, languages, availability',
)

feature_slide(
    WHITE, 'ADMIN · VIC CLIENTS',
    'High-value clients always have a matched advisor on the floor.',
    'Platinum, gold and silver clients each carry preferred languages and a list of affiliated advisors — the roster engine treats an unstaffed VIC appointment as a scoring penalty, not an afterthought.',
    ['Tiered client profiles with preferred-language matching',
     'Advisor assignment independent of the roster-generation run',
     'Upcoming appointments tracked per client, with status'],
    A.format('admin-vic.png'), image_right=False,
    caption='VIC client card — tier, advisors, upcoming appointment',
)

feature_slide(
    WHITE, 'ADMIN · SHIFTS',
    'Every boutique defines its own shifts and staffing minimums.',
    'Shift times, sort order and validity windows are configured per boutique, each with its own per-skill minimum headcount — plus a closure calendar the engine refuses to schedule against.',
    ['Independent shift definitions per boutique — no shared template lock-in',
     'Minimum (and optional maximum) headcount per skill, per shift',
     'Closure dates the engine automatically respects'],
    A.format('admin-shifts.png'), image_right=True,
    caption='Shift definitions with per-skill minimum staffing',
)

feature_slide(
    WHITE, 'ADMIN · RULES',
    'Tune the engine without touching a line of code.',
    'Eight constraint rules — hours, fatigue, compliance, coverage and availability — can each be switched between warning and hard-block, independently per boutique.',
    ['Hard-block vs. warning severity, set per rule',
     'Target headcount, max hours, rest hours and VIC boost all editable',
     'Scoring weights validated to sum to 1.0 before saving'],
    A.format('admin-rules.png'), image_right=False,
    caption='Eight constraint rules, engine thresholds and scoring weights — all per boutique',
)

feature_slide(
    WHITE, 'ADMIN · LEAVE',
    'Leave feeds straight into who the engine considers available.',
    'Annual, sick, parental, unpaid and public-holiday leave are tracked per staff member, tagged by source, and automatically exclude that person from generation on those dates.',
    ['Bulk "public holiday" entry for every staff member in one action',
     'HR-synced records are protected from accidental deletion',
     'Filterable by staff, leave type, source and date range'],
    A.format('admin-leave.png'), image_right=True,
    caption='Leave register — HR-synced, ad-hoc and manual records side by side',
)

# ─────────────────────────────────────────────────────────────────────────────
# APPROVER (two images: inbox + history)
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide(WHITE)
text_w = CW - IMG_MAX_W - Inches(0.45)
img_x = SW - MX - IMG_MAX_W
add_text(s, MX, Inches(0.75), text_w, Inches(0.35), 'APPROVER · INBOX & HISTORY', size=10.5, color=GOLD, bold=True)
add_text(s, MX, Inches(1.12), text_w, Inches(1.35), 'A second pair of eyes, kept structurally separate from the drafter.',
          size=21, color=NAVY, font=SERIF, line_spacing=1.15)
add_text(s, MX, Inches(2.55), text_w, Inches(1.9),
          'Approvers get a dedicated inbox of submitted rosters — score, overrides and full shift detail — with approve, reject (with a note back to the admin) and publish.',
          size=12, color=MUTED, line_spacing=1.32)
add_bullets(s, MX, Inches(4.55), text_w, Inches(2.3), [
    'Rejection returns a roster to draft with a reason, no dead ends',
    'Publishing safely replaces any previously published roster for that date',
    'Full approval history — published, rejected and archived — kept separately',
], size=11.5, gap=9)

half_h = Inches(2.85)
pic1, l1, t1, w1, h1 = add_image_fit(s, A.format('approver-inbox.png'), img_x, Inches(0.9), IMG_MAX_W, half_h, center_h=True)
add_text(s, img_x, t1 + h1 + Inches(0.06), IMG_MAX_W, Inches(0.3), 'Approval Inbox', size=9.5, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
pic2, l2, t2, w2, h2 = add_image_fit(s, A.format('approver-history.png'), img_x, Inches(0.9) + half_h + Inches(0.45), IMG_MAX_W, half_h, center_h=True)
add_text(s, img_x, t2 + h2 + Inches(0.06), IMG_MAX_W, Inches(0.3), 'Approval History — amended vs. published, kept apart', size=9.5, color=MUTED, italic=True, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# STAFF · MY SCHEDULE (portrait/mobile)
# ─────────────────────────────────────────────────────────────────────────────
feature_slide(
    WHITE, 'STAFF · MY SCHEDULE',
    'Every employee can check their own shifts without asking a manager.',
    'Once a roster is published, each staff member with a linked account sees their own upcoming shifts — shift name, VIC flag and duration — grouped by day, on a mobile-first layout built for checking on the go.',
    ['Reads only published rosters — never a draft in progress',
     'No admin or approver access required, just a linked staff record',
     'Leave requests and a team view are scoped for a later phase'],
    A.format('staff-my-schedule.png'), image_right=True, mobile=True,
    caption="My Schedule — an employee's own view, on their own phone",
)

# ─────────────────────────────────────────────────────────────────────────────
# SCENARIO: ad-hoc availability change
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide(WARM)
eyebrow(s, MX, Inches(0.42), 'BUSINESS AGILITY, IN PRACTICE')
heading(s, MX, Inches(0.74), CW, 'A staff member calls in sick — the morning of a published shift.', size=21)
add_text(s, MX, Inches(1.5), CW, Inches(0.55),
          'A real run-through, captured from the app: a rostered stylist becomes unavailable after her shift is already live, and the correction goes out through the normal review path in minutes.',
          size=11.5, color=MUTED, line_spacing=1.25)

frames = [
    ('1', 'LOG THE LEAVE', 'Admin · Leave', 'scenario-leave-added.png',
     'The admin adds an ad-hoc sick-leave record for today, on the spot — timestamped and attributed like every other leave record.'),
    ('2', 'REGENERATE', 'Admin · Rosters', 'scenario-new-draft.png',
     'Re-running Generate for the same date automatically excludes her and fills the shift from the bench — no one hunts for a replacement.'),
    ('3', 'RE-APPROVE & PUBLISH', 'Approver · History', 'scenario-history-superseded.png',
     'Submitted and published through the same approval step as any roster. The record for this date now shows both versions.'),
]
gap = Inches(0.3)
fw = (CW - gap * 2) / 3
fy = Inches(2.15)
card_h = Inches(4.0)
img_h = Inches(2.0)
for i, (num, tlabel, role, img, body) in enumerate(frames):
    fx = MX + i * (fw + gap)
    add_rect(s, fx, fy, fw, card_h, fill=WHITE, line_color=LINE, line_w=Pt(0.75), round_=True)
    circ = slide_shapes_add_oval(s, fx + Inches(0.18), fy - Inches(0.15), Inches(0.4), Inches(0.4))
    add_text(s, fx + Inches(0.18), fy - Inches(0.15), Inches(0.4), Inches(0.4), num, size=13, color=GOLD_LT,
              font=SERIF, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_image_fit(s, os.path.join(ASSETS, img), fx + Inches(0.18), fy + Inches(0.3), fw - Inches(0.36), img_h, center_h=True)
    add_text(s, fx + Inches(0.2), fy + Inches(0.3) + img_h + Inches(0.08), fw - Inches(0.4), Inches(0.22), tlabel,
              size=9.5, color=GOLD, bold=True)
    add_text(s, fx + Inches(0.2), fy + Inches(0.3) + img_h + Inches(0.32), fw - Inches(0.4), Inches(0.28), role,
              size=11.5, color=NAVY, bold=True)
    add_text(s, fx + Inches(0.2), fy + Inches(0.3) + img_h + Inches(0.62), fw - Inches(0.4), Inches(0.95), body,
              size=9.5, color=MUTED, line_spacing=1.18)

take_y = fy + card_h + Inches(0.2)
take_h = Inches(0.95)
add_rect(s, MX, take_y, CW, take_h, fill=NAVY_DEEP, round_=True)
add_text(s, MX + Inches(0.35), take_y + Inches(0.1), Inches(0.5), Inches(0.5), '✦', size=18, color=GOLD)
add_text(s, MX + Inches(0.85), take_y + Inches(0.09), CW - Inches(1.2), Inches(0.8),
          'Same day, one date, two versions, zero manual reconciliation: the corrected roster is published, the original is kept as an Amended record rather than deleted.',
          size=10.5, color=RGBColor.from_string('D8D7E8'), line_spacing=1.2)

# ─────────────────────────────────────────────────────────────────────────────
# BUDGET
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide(NAVY_DEEP)
eyebrow(s, MX, Inches(0.42), 'BUDGET PLANNING', color=GOLD)
add_text(s, MX, Inches(0.74), CW, Inches(0.8), "What each role's screen time is worth.", size=24, color=WHITE, font=SERIF)
add_text(s, MX, Inches(1.4), CW, Inches(0.4),
          'Framed the way a budget review usually wants it: who uses which portal, and what it replaces.',
          size=11.5, color=RGBColor.from_string('B9B8C9'))

rows = [
    ('Admin', 'Rosters, Staff, VIC, Shifts, Rules, Leave', 'Store or regional operations admin',
     'Manual spreadsheet rostering, side-channel leave tracking',
     'One generation run replaces hours of manual shift-building per boutique per week'),
    ('Approver', 'Inbox, History', 'Boutique or regional manager',
     'Ad-hoc email/print sign-off with no audit trail',
     'A single queue with a permanent, searchable approval record'),
    ('Staff', 'My Schedule', 'Every floor employee',
     '"What\'s my shift?" messages to the manager',
     "Removes a constant low-value interruption from every admin's day"),
    ('Reader', 'Published Rosters — planned', 'Shared/back-of-house display',
     'Printed roster taped to a wall',
     'Always-current view with no reprint cycle once delivered'),
]
tbl_top = Inches(2.0)
tbl_left = MX
tbl_w = CW
tbl_h = Inches(3.15)
rows_n = len(rows) + 1
gtable = s.shapes.add_table(rows_n, 4, tbl_left, tbl_top, tbl_w, tbl_h).table
gtable.columns[0].width = Inches(2.6)
gtable.columns[1].width = Inches(2.9)
gtable.columns[2].width = Inches(3.4)
gtable.columns[3].width = Inches(2.93)
headers = ['Portal', 'Typical user', 'Replaces / reduces', 'Value']
for c, htext in enumerate(headers):
    cell = gtable.cell(0, c)
    cell.fill.solid()
    cell.fill.fore_color.rgb = NAVY_DEEP
    tf = cell.text_frame
    tf.paragraphs[0].text = htext.upper()
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = GOLD
    tf.paragraphs[0].font.name = SANS
for r, (role, sub, user, replaces, value) in enumerate(rows, start=1):
    vals = [f'{role}\n{sub}', user, replaces, value]
    for c, val in enumerate(vals):
        cell = gtable.cell(r, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY_DEEP
        cell.margin_top = Pt(6)
        cell.margin_bottom = Pt(6)
        tf = cell.text_frame
        tf.word_wrap = True
        lines = val.split('\n')
        for li, line in enumerate(lines):
            p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
            p.line_spacing = 1.15
            r_ = p.add_run()
            r_.text = line
            r_.font.size = Pt(10 if (c == 0 and li == 0) else 9)
            r_.font.name = SANS
            if c == 0 and li == 0:
                r_.font.bold = True
                r_.font.color.rgb = WHITE
            else:
                r_.font.color.rgb = RGBColor.from_string('9C9BB5') if (c == 0 and li == 1) else RGBColor.from_string('C7C6DA')

note_y = Inches(5.4)
add_rect(s, MX, note_y, CW, Inches(1.05), fill=RGBColor.from_string('1C2456'), round_=True)
add_text(s, MX + Inches(0.3), note_y + Inches(0.12), CW - Inches(0.6), Inches(0.3), 'COST OF SCALING TO MORE BOUTIQUES',
          size=10.5, color=GOLD_LT, bold=True)
add_text(s, MX + Inches(0.3), note_y + Inches(0.42), CW - Inches(0.6), Inches(0.6),
          'Staff, VIC clients, shifts, rules and scoring weights are configured per boutique inside the same deployment. Adding a store is a data-entry and configuration exercise, not a new engineering project.',
          size=10, color=RGBColor.from_string('B9B8C9'), line_spacing=1.2)

chip_y = Inches(6.65)
chips = ['No backend server to run', 'No compute costs for the scheduling engine', 'One data model across every boutique', 'Auto-deploy on every change']
cx = MX
for chip in chips:
    cw_ = Inches(0.35 + 0.075 * len(chip))
    add_rect(s, cx, chip_y, cw_, Inches(0.42), fill=RGBColor.from_string('242C5C'), round_=True)
    add_text(s, cx + Inches(0.12), chip_y + Inches(0.07), cw_ - Inches(0.24), Inches(0.3), chip, size=9.5,
              color=RGBColor.from_string('C7C6DA'), align=PP_ALIGN.CENTER)
    cx += cw_ + Inches(0.18)

# ─────────────────────────────────────────────────────────────────────────────
# ROADMAP
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide(WARM)
eyebrow(s, MX, Inches(0.5), "WHAT'S NEXT")
heading(s, MX, Inches(0.85), CW, 'Scoped, not hidden.', size=26)
add_text(s, MX, Inches(1.65), CW, Inches(0.5),
          'A few screens are intentionally stubbed for a later phase rather than rushed — worth knowing for planning purposes.',
          size=12.5, color=MUTED)

road = [
    ('Staff leave requests', 'Let employees submit their own unavailability instead of going through an admin.'),
    ('Team view for staff', 'Let staff see the published roster for their whole boutique, not just their own shifts.'),
    ('Reader / published-roster browser', 'Read-only roster browsing with PDF/CSV export for shared displays and print.'),
]
gap = Inches(0.3)
rw = (CW - gap * 2) / 3
ry = Inches(2.5)
rh = Inches(2.6)
for i, (title, body) in enumerate(road):
    rx = MX + i * (rw + gap)
    add_rect(s, rx, ry, rw, rh, fill=WHITE, line_color=LINE, line_w=Pt(0.75), round_=True)
    add_pill(s, rx + Inches(0.25), ry + Inches(0.25), 'Phase 4', AMBER_TX, AMBER_BG, size=9)
    add_text(s, rx + Inches(0.25), ry + Inches(0.75), rw - Inches(0.5), Inches(0.8), title, size=14.5, color=NAVY,
              bold=True, line_spacing=1.15)
    add_text(s, rx + Inches(0.25), ry + Inches(1.55), rw - Inches(0.5), Inches(0.9), body, size=11, color=MUTED, line_spacing=1.3)

# ─────────────────────────────────────────────────────────────────────────────
# CTA
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide(NAVY_DEEP)
add_text(s, Inches(2), Inches(2.7), SW - Inches(4), Inches(0.6), '✦', size=30, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s, Inches(1.5), Inches(3.3), SW - Inches(3), Inches(1.3),
          'Four roles. One system. Built to add the\nnext boutique, not rebuild for it.',
          size=27, color=WHITE, font=SERIF, align=PP_ALIGN.CENTER, line_spacing=1.2)
add_text(s, Inches(2.5), Inches(4.75), SW - Inches(5), Inches(0.5), 'Product overview for internal budget review',
          size=13, color=RGBColor.from_string('9C9BB5'), align=PP_ALIGN.CENTER)

prs.save(OUT)
print('Saved full deck:', OUT, '-', len(prs.slides._sldIdLst), 'slides')
